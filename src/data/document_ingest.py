"""Teacher upload path: turn an uploaded worksheet/document into a course_content row.

Text-based formats (.docx, .pdf, .pptx, .txt, .md) are parsed deterministically here.
Images go through a separate LLM-based path (agents/vision_agent.py) since extracting
their content isn't a deterministic operation - see app.py's upload handler, which calls
`store_course_content` directly with the vision agent's transcription.

This is the "how does real ManageBac/Teams content get in" answer for a no-budget
course project: a teacher drops a file in via the Streamlit sidebar instead of a
live platform integration.

`store_course_content` is also the single choke point both upload paths (text here,
images in app.py) funnel through, so it's where the row gets indexed into the vector
search store (capabilities/retrieval/vector_store.py) right after the SQLite insert -
one hook keeps both upload paths in sync rather than duplicating the call at each
call site.
"""

from __future__ import annotations

import sqlite3
from datetime import datetime, timezone
from pathlib import Path
from typing import BinaryIO

from docx import Document
from pptx import Presentation
from pypdf import PdfReader

from capabilities.retrieval import vector_store
from capabilities.retrieval.vector_store import Collection

SUPPORTED_DOCUMENT_EXTENSIONS = {".docx", ".pdf", ".pptx", ".txt", ".md"}


def extract_text_from_docx(file: BinaryIO | str) -> str:
    """Return the visible paragraph text of a .docx file, one paragraph per line."""
    document = Document(file)
    paragraphs = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def extract_text_from_pdf(file: BinaryIO | str) -> str:
    """Return the extracted text of a .pdf file, one page per paragraph break."""
    reader = PdfReader(file)
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(page.strip() for page in pages if page.strip())


def extract_text_from_pptx(file: BinaryIO | str) -> str:
    """Return the visible text of a .pptx file, one slide per paragraph break."""
    presentation = Presentation(file)
    slides_text = []
    for slide in presentation.slides:
        lines = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if lines:
            slides_text.append("\n".join(lines))
    return "\n\n".join(slides_text)


def extract_text_from_plain_text(file: BinaryIO | str) -> str:
    """Return the content of a .txt/.md file as-is."""
    if isinstance(file, (str, Path)):
        return Path(file).read_text(encoding="utf-8").strip()
    raw = file.read()
    if isinstance(raw, bytes):
        raw = raw.decode("utf-8")
    return raw.strip()


_EXTRACTORS = {
    ".docx": extract_text_from_docx,
    ".pdf": extract_text_from_pdf,
    ".pptx": extract_text_from_pptx,
    ".txt": extract_text_from_plain_text,
    ".md": extract_text_from_plain_text,
}


def extract_text(file: BinaryIO | str, filename: str) -> str:
    """Dispatch to the right extractor for `filename`'s extension."""
    suffix = Path(filename).suffix.lower()
    extractor = _EXTRACTORS.get(suffix)
    if extractor is None:
        raise ValueError(
            f"Unsupported file type {suffix!r} for {filename!r}; "
            f"expected one of {sorted(SUPPORTED_DOCUMENT_EXTENSIONS)}"
        )
    return extractor(file)


def store_course_content(
    conn: sqlite3.Connection,
    *,
    subject: str,
    topic: str,
    content: str,
    source_name: str,
    collection: Collection,
    uploader_role: str = "Teacher",
    uploader_name: str = "",
) -> int:
    """Insert already-extracted text as a course_content row, and index it into the
    vector search store so it's immediately searchable. Returns the new row's id.

    `uploader_role`/`uploader_name` record who added it (e.g. "Teacher"/"Ms. Smith")
    for the course library panel; `created_at` is stamped as of this call so the
    panel can order newest-first.
    """
    if not content:
        raise ValueError(f"No content to store for {source_name!r}")

    created_at = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")
    cursor = conn.execute(
        "INSERT INTO course_content (subject, topic, content, source, owner, uploaded_by, created_at) "
        "VALUES (?, ?, ?, ?, ?, ?, ?)",
        (subject, topic, content, source_name, uploader_role, uploader_name, created_at),
    )
    conn.commit()
    row_id = cursor.lastrowid
    vector_store.index_row(
        collection,
        {"id": row_id, "subject": subject, "topic": topic, "content": content, "source": source_name},
    )
    return row_id


def ingest_document(
    conn: sqlite3.Connection,
    *,
    file: BinaryIO | str,
    subject: str,
    topic: str,
    source_name: str,
    collection: Collection,
    uploader_role: str = "Teacher",
    uploader_name: str = "",
) -> int:
    """Parse an uploaded document (.docx/.pdf/.pptx/.txt/.md) and store it as a
    course_content row. Returns the new row's id.
    """
    text = extract_text(file, source_name)
    return store_course_content(
        conn,
        subject=subject,
        topic=topic,
        content=text,
        source_name=source_name,
        collection=collection,
        uploader_role=uploader_role,
        uploader_name=uploader_name,
    )


def list_course_content(conn: sqlite3.Connection) -> list[dict]:
    """id, filename, subject, owner, uploaded_by, created_at for every row - powers
    the library panel, newest upload first.

    `source` is aliased to `filename` (the UI-facing name for "uploaded file
    name, or 'seed'"). No `content` body - the panel only lists files. Ties on
    `created_at` (e.g. seed rows, which all share `db.SEED_UPLOAD_TIMESTAMP`)
    break on `id` descending, so ordering is always deterministic.
    """
    rows = conn.execute(
        "SELECT id, source AS filename, subject, owner, uploaded_by, created_at "
        "FROM course_content ORDER BY created_at DESC, id DESC"
    ).fetchall()
    return [dict(row) for row in rows]


def get_course_content_by_id(conn: sqlite3.Connection, content_id: int) -> dict | None:
    """Full row (id, filename, subject, owner, content) for folding into a
    prompt, or None if content_id doesn't exist (defensive - a stale
    session_state id)."""
    row = conn.execute(
        "SELECT id, source AS filename, subject, owner, content "
        "FROM course_content WHERE id = ?",
        (content_id,),
    ).fetchone()
    return dict(row) if row else None
