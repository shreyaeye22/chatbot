from __future__ import annotations

import io

import pytest
from docx import Document
from pptx import Presentation
from reportlab.pdfgen import canvas

from data.db import get_connection, init_db
from data.document_ingest import (
    extract_text,
    extract_text_from_docx,
    extract_text_from_pdf,
    extract_text_from_plain_text,
    extract_text_from_pptx,
    ingest_document,
    store_course_content,
)


@pytest.fixture
def db_conn(tmp_path):
    db_path = tmp_path / "app.db"
    init_db(db_path)
    conn = get_connection(db_path)
    yield conn
    conn.close()


def test_extract_text_from_docx(tmp_path):
    path = tmp_path / "worksheet.docx"
    doc = Document()
    doc.add_paragraph("Comparative and superlative worksheet")
    doc.add_paragraph("Complete exercises 1 to 10.")
    doc.save(path)

    text = extract_text_from_docx(str(path))

    assert "Comparative and superlative worksheet" in text
    assert "Complete exercises 1 to 10." in text


def test_extract_text_from_pdf(tmp_path):
    path = tmp_path / "notes.pdf"
    c = canvas.Canvas(str(path))
    c.drawString(72, 700, "Population pyramids show age and gender structure.")
    c.save()

    text = extract_text_from_pdf(str(path))

    assert "Population pyramids show age and gender structure." in text


def test_extract_text_from_pptx(tmp_path):
    path = tmp_path / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Cell structure"
    slide.placeholders[1].text = "Plant cells have a cell wall and chloroplasts."
    presentation.save(path)

    text = extract_text_from_pptx(str(path))

    assert "Cell structure" in text
    assert "Plant cells have a cell wall and chloroplasts." in text


def test_extract_text_from_plain_text_path(tmp_path):
    path = tmp_path / "notes.txt"
    path.write_text("Newton's second law: F = m * a", encoding="utf-8")

    assert extract_text_from_plain_text(str(path)) == "Newton's second law: F = m * a"


def test_extract_text_from_plain_text_file_like():
    file = io.BytesIO("Balancing equations conserves mass.".encode("utf-8"))

    assert extract_text_from_plain_text(file) == "Balancing equations conserves mass."


def test_extract_text_dispatches_docx_by_extension(tmp_path):
    path = tmp_path / "worksheet.docx"
    doc = Document()
    doc.add_paragraph("dispatch check")
    doc.save(path)

    assert extract_text(str(path), "worksheet.docx") == extract_text_from_docx(str(path))


def test_extract_text_dispatches_txt_by_extension_case_insensitively(tmp_path):
    path = tmp_path / "notes.TXT"
    path.write_text("dispatch check", encoding="utf-8")

    assert extract_text(str(path), "notes.TXT") == "dispatch check"


def test_extract_text_rejects_unsupported_extension():
    with pytest.raises(ValueError):
        extract_text("fake-file-handle", "notes.xlsx")


def test_store_course_content_inserts_a_row(db_conn):
    row_id = store_course_content(
        db_conn, subject="math", topic="algebra", content="some notes", source_name="upload.txt"
    )

    row = db_conn.execute(
        "SELECT subject, topic, content, source FROM course_content WHERE id = ?", (row_id,)
    ).fetchone()

    assert row["subject"] == "math"
    assert row["topic"] == "algebra"
    assert row["content"] == "some notes"
    assert row["source"] == "upload.txt"


def test_store_course_content_rejects_empty_content(db_conn):
    with pytest.raises(ValueError):
        store_course_content(
            db_conn, subject="math", topic="algebra", content="", source_name="empty.txt"
        )


def test_ingest_document_end_to_end(tmp_path, db_conn):
    path = tmp_path / "notes.txt"
    path.write_text("The mitochondria is the powerhouse of the cell.", encoding="utf-8")

    row_id = ingest_document(
        db_conn, file=str(path), subject="biology", topic="cells", source_name="notes.txt"
    )

    row = db_conn.execute("SELECT content FROM course_content WHERE id = ?", (row_id,)).fetchone()
    assert "mitochondria" in row["content"]
